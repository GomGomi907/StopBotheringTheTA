import logging
import re
from pathlib import Path
from typing import Dict, Any, Optional

# import fitz  # Removed unused dependency
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

class MetadataExtractor:
    def __init__(self):
        pass

    def extract_text_from_file(self, file_path: Path) -> str:
        """파일 경로에서 텍스트를 추출 (PDF, PPTX, DOCX 지원)"""
        if not file_path.exists():
            return ""
        
        ext = file_path.suffix.lower()
        text = ""
        
        # 파일 크기 확인 (0바이트 무시)
        if file_path.stat().st_size == 0:
            return ""

        try:
            if ext == ".pdf":
                try:
                    reader = PdfReader(str(file_path))
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                except Exception:
                    pass # 암호화된 PDF 등
            elif ext == ".pptx":
                from pptx.exc import PackageNotFoundError
                try:
                    prs = Presentation(str(file_path))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                except PackageNotFoundError:
                    # 다운로드 실패로 HTML 에러 페이지가 저장된 경우일 수 있음
                    pass
            elif ext == ".docx":
                from docx.opc.exceptions import PackageNotFoundError
                try:
                    doc = Document(str(file_path))
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                except PackageNotFoundError:
                    pass
            else:
                # 텍스트 파일 등 시도
                try:
                    text = file_path.read_text(encoding="utf-8", errors="ignore")
                except:
                    pass
        except Exception as e:
            # 기타 에러는 로그만 남기고 무시 (프로세스 중단 방지)
            logger.debug(f"텍스트 추출 건너뜀 ({file_path.name}): {e}")
            
        return text.strip()

    def summarize_record(self, record: Dict[str, Any]) -> Dict[str, Any]:
        """레코드(공지, 과제 등)에서 요약 메타데이터 추출"""
        return self.summarize_item(record.get("payload", {}), record.get("category"), record.get("title"))

    def summarize_item(self, payload: Dict[str, Any], category: str, default_title: str = "No Title") -> Dict[str, Any]:
        """단일 아이템(Payload Dict)에서 메타데이터 추출"""
        # 아이템 내부에 title/name이 있으면 우선 사용
        title = payload.get("title") or payload.get("name") or default_title
        
        meta = {
            "title": title,
            "category": category,
            "url": payload.get("html_url") or payload.get("url"), # Canvas Item URL
            "content_summary": "",
            "date": "",
        }

        # 날짜/내용 추출 로직
        if isinstance(payload, dict):
            # Canvas Content
            # Pages API는 'body'에 내용이 있음
            html = payload.get("body") or payload.get("message") or payload.get("description") or payload.get("html")
            if html:
                try:
                    soup = BeautifulSoup(html, "html.parser")
                    text = soup.get_text(separator=" ", strip=True)
                    meta["content_summary"] = text[:1000]  # 요약 길이 증가 (500 -> 1000)
                except:
                    meta["content_summary"] = str(html)[:500]

            # 날짜 (우선순위: due_at > posted_at > created_at)
            meta["date"] = payload.get("due_at") or payload.get("posted_at") or payload.get("created_at")
            
            # [Added] Context-Aware Logic을 위한 원본 날짜 보존
            meta["due_at"] = payload.get("due_at")
            meta["posted_at"] = payload.get("posted_at") or payload.get("created_at")
        
        return meta
